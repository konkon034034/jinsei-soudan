from pptx import Presentation
from pptx.util import Inches, Pt
import os

# Define the content
slides_content = [
    {
        "title": "第5位：きなこ棒 (Kinako Bo)",
        "image": r"C:\Users\rentsu-windows\.gemini\antigravity\brain\4ddab301-7bf4-4a03-b3e3-4c23851f5878\kinako_bo_retro_1765688680720.png",
        "description": "きな粉と水飴を練り合わせて作られた素朴な駄菓子。\n独特の甘さとねっとりとした食感が特徴で、当たりの爪楊枝が出るともう一本もらえるワクワク感も人気の秘密でした。"
    },
    {
        "title": "第4位：ココアシガレット (Cocoa Cigarette)",
        "image": r"C:\Users\rentsu-windows\.gemini\antigravity\brain\4ddab301-7bf4-4a03-b3e3-4c23851f5878\cocoa_cigarette_retro_1765688663132.png",
        "description": "タバコのような形状をしたココア味の菓子。\n子供たちが大人の真似をして「プカプカ」とふかす真似をして楽しむ姿は、昭和の路地裏の定番風景でした。"
    },
    {
        "title": "第3位：よっちゃんイカ (Yocchan Ika)",
        "image": r"C:\Users\rentsu-windows\.gemini\antigravity\brain\4ddab301-7bf4-4a03-b3e3-4c23851f5878\yocchan_ika_retro_1765688644345.png",
        "description": "イカを加工した酢漬けの駄菓子。\nその強烈な酸味と独特の食感は一度食べると病みつきに。当たりくじ付きのパッケージも子供心をくすぐりました。"
    },
    {
        "title": "第2位：キャベツ太郎 (Cabbage Taro)",
        "image": r"C:\Users\rentsu-windows\.gemini\antigravity\brain\4ddab301-7bf4-4a03-b3e3-4c23851f5878\cabbage_taro_retro_1765688627380.png",
        "description": "丸い形とサクサクとした食感が特徴のスナック菓子。\nキャベツが入っているわけではありませんが、その丸い形が愛らしく、ソース味が後を引く美味しさです。"
    },
    {
        "title": "第1位：うまい棒 (Umai-bo)",
        "image": r"C:\Users\rentsu-windows\.gemini\antigravity\brain\4ddab301-7bf4-4a03-b3e3-4c23851f5878\umai_bo_retro_1765688610757.png",
        "description": "1979年発売の王道駄菓子。\n1本10円（当時）という安さと、コーンポタージュ、めんたい、チーズなど豊富なフレーバーで、不動の人気No.1を誇ります。"
    }
]

def create_presentation():
    prs = Presentation()

    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "昭和の駄菓子屋 人気TOP5"
    subtitle.text = "懐かしの味と記憶を振り返る"

    # Content Slides
    for item in slides_content:
        # Use a blank layout for custom positioning
        blank_slide_layout = prs.slide_layouts[6] 
        slide = prs.slides.add_slide(blank_slide_layout)

        # Add Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = item["title"]
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = Pt(32)
        title_paragraph.font.bold = True

        # Add Image
        if os.path.exists(item["image"]):
            slide.shapes.add_picture(item["image"], Inches(0.5), Inches(1.8), height=Inches(3.5))
        else:
            print(f"Warning: Image not found at {item['image']}")

        # Add Description
        tb = slide.shapes.add_textbox(Inches(5.5), Inches(1.8), Inches(4), Inches(3.5))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.text = item["description"]
        tf.paragraphs[0].font.size = Pt(18)

    output_path = "showa_dagashi_top5.pptx"
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    create_presentation()
